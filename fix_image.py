from pptx import Presentation

prs = Presentation("DrishtiAI_SIH2026_Pitch_SIH_Enhanced.pptx")
slide_3 = prs.slides[2]

target_shape = None
for s in slide_3.shapes:
    if getattr(s, "left", 0) and abs(s.left - 8055864) < 10000 and abs(s.top - 1051560) < 10000:
        target_shape = s
        break

if target_shape:
    left, top, w, h = target_shape.left, target_shape.top, target_shape.width, target_shape.height
    
    # Remove the old shape (which was a generic AUTO_SHAPE with picture fill)
    sp = target_shape._element
    sp.getparent().remove(sp)
    
    # Add the new picture exactly where the old shape was
    slide_3.shapes.add_picture("pipeline.png", left, top, width=w, height=h)
    print("Replaced patient scan shape with pipeline.png!")
else:
    print("Could not find target shape!")

prs.save("DrishtiAI_SIH2026_Pitch_Final_Enhanced.pptx")
print("Saved to DrishtiAI_SIH2026_Pitch_Final_Enhanced.pptx")
